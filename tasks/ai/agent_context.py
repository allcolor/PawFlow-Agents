"""AgentLoopTask mixin — AgentContext methods

Auto-extracted from tasks/ai/agent_loop.py.
All methods access self (AgentLoopTask instance).
"""
import logging
import os
from typing import Dict, Any, List, Optional


from core import FlowFile
from core.llm_client import (
    LLMMessage,
)
from tasks.ai._agentctx_base import (  # noqa: F401
    _PACState, _find_agent_md, _agent_md_cache, _AGENT_MD_TTL)
from tasks.ai.agent_tool_config import AgentToolConfigMixin
from tasks.ai.agent_tool_exec import AgentToolExecMixin
from tasks.ai._agentctx_p1 import _PACPhase1Mixin
from tasks.ai._agentctx_p2 import _PACPhase2Mixin
from tasks.ai._agentctx_p3 import _PACPhase3Mixin

logger = logging.getLogger(__name__)


class _PawFlowMarkItDownVisionClient:
    """Minimal OpenAI-style vision client backed by a PawFlow LLM service."""

    def __init__(self, service, user_id: str = "", conversation_id: str = ""):
        self._service = service
        self._user_id = user_id
        self._conversation_id = conversation_id or "markitdown-ocr"
        self.chat = self._Chat(self)

    class _Chat:
        def __init__(self, outer):
            self.completions = outer._Completions(outer)

    class _Completions:
        def __init__(self, outer):
            self._outer = outer

        def create(self, model="", messages=None, **kwargs):
            from types import SimpleNamespace
            from core.llm_client import LLMMessage
            prompt_messages = []
            for msg in messages or []:
                role = str(msg.get("role") or "user")
                content = msg.get("content", "")
                prompt_messages.append(LLMMessage(
                    role=role,
                    content=content,
                    conversation_id=self._outer._conversation_id,
                ))
            response = self._outer._service.complete(
                prompt_messages,
                model=model or None,
                temperature=float(kwargs.get("temperature", 0) or 0),
                max_tokens=int(kwargs.get("max_tokens", 1000) or 1000),
                call_user_id=self._outer._user_id,
                call_conversation_id=self._outer._conversation_id,
            )
            return SimpleNamespace(choices=[SimpleNamespace(
                message=SimpleNamespace(content=response.content or ""))])


class AgentContextMixin(AgentToolConfigMixin, AgentToolExecMixin,
                        _PACPhase1Mixin, _PACPhase2Mixin, _PACPhase3Mixin):
    """Context preparation + user content building."""

    def _pac_cfg(self, st, key, default):
        """Agent overrides service, service overrides default.
        None or empty string = not set. 0 IS a valid override."""
        v = self.config.get(key)
        if v is not None and v != "":
            return v
        v = st._svc_cfg.get(key)
        if v is not None and v != "":
            return v
        return default

    def _prepare_agent_context(self, flowfile: FlowFile, *,
                               preloaded_messages: Optional[List[Dict]] = None,
                               preloaded_conversation_id: str = "",
                               independent_context: bool = False):
        """Extract common context from flowfile and config for both sync and streaming modes.

        Args:
            flowfile: The FlowFile with request data.
            preloaded_messages: If set, use these raw message dicts instead of
                loading from ConversationStore. Used by the poller for task
                sub-conversations that have their own isolated message store.
        """
        st = _PACState()
        st.flowfile = flowfile
        st.preloaded_messages = preloaded_messages
        st.preloaded_conversation_id = preloaded_conversation_id
        st.independent_context = independent_context
        self._pac_p1(st)
        self._pac_p2(st)
        self._pac_p3(st)
        return {
            "client": st.client, "registry": st.registry, "tool_defs": st.tool_defs,
            "messages": st.messages, "model": st.model_name,
            "_turn_mode": st._turn_mode,
            "_identity_suffix": st._identity_suffix,
            "temperature": st.temperature, "max_tokens": st.max_tokens,
            "max_iterations": st.max_iterations,
            "max_consecutive_tool_calls": st.max_consecutive_tool_calls,
            "thinking_budget": st.thinking_budget,
            "max_rounds": int(st._cfg("max_rounds", 1)),
            "use_conv_store": st.use_conv_store, "conv_ttl": st.conv_ttl,
            "conv_attr": st.conv_attr, "conversation_id": st.conversation_id,
            "user_id": st.user_id,
            "_base_message_count": st.base_message_count,
            "max_context_size": int(
                # Per-agent: use service max_tokens (= context window size)
                st._resolved_max_ctx
            ),
            "configured_context_size": int(st._configured_max_ctx or 0),
            "real_context_size": int(st._real_max_ctx or 0),
            "context_keep_recent": int(st._cfg("context_keep_recent", 6)),
            "chars_per_token": float(
                (getattr(st.resolved_svc, 'config', {}) or {}).get("chars_per_token", 0)
                or self.config.get("chars_per_token", 0)
            ),
            "channel": st.channel,
            "request_msg_id": st.flowfile.get_attribute("agent.request_msg_id") or "",
            "active_agent_name": st._active_agent_name,  # MUST be non-empty — see _ensure_active_agent
            "active_llm_service": st._active_llm_service,
            "title_llm_service": self._resolve_service_param(
                "title_llm_service", st.user_id, st.conversation_id),
            "resolved_svc": st.resolved_svc,
            "max_budget_usd": st._max_budget,
            "summarizer": self._get_summarizer_client(st.user_id, conversation_id=st.conversation_id),  # (client, max_ctx, svc_id)
            "sub_executor": st.sub_executor,
            "_target_agent": st._target_agent,
            "_context_diverged": st._context_diverged,
            "_materialize_pawflow_initial_context": bool(st._uses_pawflow_initial),
            "_pawflow_initial_context_source": st._cold_cli_initial_source,
            "_nicknames": st._nicknames if st.conversation_id else {},
            "_is_cli_provider": st._is_cli_provider,
            "_cli_has_session": st._cli_has_session,
            "_is_claude_code": st._is_claude_code or st._is_claude_code_interactive,
            "_claude_has_session": st._claude_has_session,
            "_agent_md_content": st._agent_md_content,
            "_provider_system_prompt": st._provider_system_prompt,
            "_datetime_str": st._datetime_str,
        }

    def _auto_compact_messages(self, messages: List[LLMMessage],
                               conversation_id: str, agent_name: str,
                               user_id: str,
                               max_context: int = 200000,
                               compact_instructions: str = "",
                               independent_context: bool = False) -> List[LLMMessage]:
        """Auto-compact if the context is past the service trigger threshold.

        Delegates to _compact which uses its own trigger_fraction (default
        0.9: only fires once real-token usage crosses 90%) and enforces
        the target_fraction hard cap (default 0.25) on its output.
        """
        _sc, _sc_max, _sc_svc = self._get_summarizer_client(user_id, conversation_id=conversation_id)
        if not _sc:
            raise RuntimeError(
                "No summarizer_service configured. Cannot compact context. "
                "Set summarizer_service in agent or flow config.")
        return self._compact(
            messages, _sc, max_context,
            conversation_id=conversation_id,
            agent_name=agent_name,
            chars_per_token=0,
            compact_instructions=compact_instructions,
            user_id=user_id,
            independent_context=independent_context,
        )

    # ── Context operation pause/resume ─────────────────────────────────


    def _build_user_content(self, text: str, attachments: List[Dict], conversation_id: str = "", user_id: str = "") -> Any:
        """Build user message content from text and optional attachments.

        If no attachments, returns plain str.
        If attachments exist, returns multi-part list for vision/document support.

        Attachment format from client:
            {"filename": "photo.png", "mime_type": "image/png", "data": "base64..."}
            {"filename": "doc.pdf", "mime_type": "application/pdf", "data": "base64..."}
        """
        if not attachments:
            return text

        import base64

        _IMAGE_TYPES = {"image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp"}
        _TEXT_TYPES = {
            "text/plain", "text/html", "text/markdown", "text/csv",
            "application/json", "application/xml",
        }
        _CONVERTIBLE_TYPES = {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
            "application/vnd.oasis.opendocument.text",  # .odt
            "application/vnd.oasis.opendocument.spreadsheet",  # .ods
            "application/msword",  # .doc (old)
            "application/vnd.ms-excel",  # .xls (old)
            "application/rtf",  # .rtf
            "application/epub+zip",  # .epub
        }
        _CONVERTIBLE_EXTS = {
            ".docx", ".xlsx", ".pptx", ".odt", ".ods",
            ".doc", ".xls", ".rtf", ".epub",
        }

        parts: List[Dict[str, Any]] = []

        # Add text first
        if text.strip():
            parts.append({"type": "text", "text": text})

        # Attachments re-materialized into the conversation FileStore get a
        # finite TTL (configurable; 24h default) instead of being permanent.
        # An attachment carrying a file_id reuses the pre-uploaded copy and
        # its TTL, so only the inline-base64 path (no prior upload) needs this.
        from core.file_ttl import resolve_ttl_seconds
        _attach_ttl = resolve_ttl_seconds(
            conversation_id=conversation_id or "",
            conv_keys=("attachment_ttl_seconds", "webchat_upload_ttl_seconds"),
            env_key="PAWFLOW_ATTACHMENT_TTL_SECONDS",
            default=86400,
        )

        for att in attachments:
            mime = att.get("mime_type", "application/octet-stream")
            filename = att.get("filename", "file")
            data_b64 = att.get("data", "")
            att_fid = att.get("file_id", "")

            # Resolve raw bytes: either from pre-uploaded file_id or inline base64
            from core.file_store import FileStore
            _fs = FileStore.instance()
            if att_fid:
                _result = _fs.get(att_fid, user_id=user_id)
                if _result:
                    _, raw, _ = _result
                else:
                    parts.append({"type": "text", "text": f"[Attached file: {filename} — upload expired]"})
                    continue
            elif data_b64:
                raw = base64.b64decode(data_b64)
            else:
                parts.append({"type": "text", "text": f"[Attached file: {filename} — no data]"})
                continue

            if mime in _IMAGE_TYPES:
                import time as _time
                # Downscale proactively to the vision ceiling so the stored
                # copy every downstream path uses (image_ref -> provider
                # materialisation -> the agent's own read/see) is already
                # within limits. Without this, a full-res phone screenshot is
                # rejected at read time ("exceeds 2000x2000") and the agent has
                # to resize it by hand. A pre-uploaded file_id is resized too,
                # but only when oversized (the helper is a no-op when it fits).
                from core.image_resize import resize_image_for_vision
                original_raw = raw
                original_mime = mime
                raw, mime = resize_image_for_vision(raw, mime)
                _ext = "jpg" if mime == "image/jpeg" else (
                    filename.rsplit('.', 1)[-1] if '.' in filename else 'png')
                _img_fname = f"image_{int(_time.time())}_{len(parts)}.{_ext}"
                # Re-store under attachment category. A reused file_id is
                # re-stored only when the resize actually changed the bytes;
                # an unchanged image keeps its original file_id.
                if att_fid and raw is original_raw and mime == original_mime:
                    _img_fid = att_fid
                else:
                    _img_fid = _fs.store(
                        _img_fname, raw, mime,
                        user_id=user_id,
                        conversation_id=conversation_id or "",
                        ttl=_attach_ttl,
                        category="attachment")
                logger.info("Attachment image: %s (%d bytes, %s) -> %s",
                            filename, len(raw), mime, _img_fid)
                parts.append({
                    "type": "image_ref",
                    "file_id": _img_fid,
                    "filename": _img_fname if not att_fid else filename,
                    "mime_type": mime,
                    "size": len(raw),
                })
            else:
                try:
                    _fid = att_fid or _fs.store(
                        filename, raw, mime,
                        user_id=user_id,
                        conversation_id=conversation_id or "",
                        ttl=_attach_ttl,
                        category="attachment")
                    logger.info("Attachment stored: %s (%s, %d bytes) -> %s",
                                filename, mime, len(raw), _fid)
                    parts.append({
                        "type": "file_ref",
                        "file_id": _fid,
                        "filename": filename,
                        "mime_type": mime,
                        "size": len(raw),
                    })
                    extracted = self._extract_attachment_markdown(
                        raw, filename, mime, conversation_id, user_id)
                    if extracted:
                        max_chars = int(
                            self.config.get("attachment_markdown_max_chars", 30000)
                            or 30000)
                        if max_chars > 0 and len(extracted) > max_chars:
                            extracted = extracted[:max_chars].rstrip() + "\n...[truncated]"
                        parts.append({
                            "type": "text",
                            "text": f"[Extracted Markdown from {filename}]\n{extracted}",
                        })
                except Exception:
                    parts.append({
                        "type": "text",
                        "text": f"[Attached file: {filename} ({mime}) — binary content, not convertible]",
                    })

        return parts if len(parts) > 1 or any(p["type"] != "text" for p in parts) else (parts[0]["text"] if parts else text)


    def _extract_attachment_markdown(self, raw: bytes, filename: str, mime: str,
                                     conversation_id: str = "",
                                     user_id: str = "") -> str:
        """Extract bounded Markdown/text for LLM context while preserving file_ref."""
        mime = (mime or "application/octet-stream").split(";", 1)[0].strip().lower()
        lower = (filename or "").lower()
        ext = lower.rsplit(".", 1)[-1] if "." in lower else ""
        text_mimes = {
            "text/plain", "text/html", "text/markdown", "text/csv",
            "application/json", "application/xml", "text/xml",
        }
        markdown_exts = {
            "pdf", "docx", "doc", "xlsx", "xls", "pptx", "odt",
            "ods", "rtf", "epub", "html", "htm", "csv", "json",
            "xml", "txt", "md", "zip",
        }
        if mime in text_mimes:
            return raw.decode("utf-8", errors="replace").strip()
        if ext not in markdown_exts and not (mime == "application/pdf" or mime.startswith("application/vnd.")):
            return ""

        converted = self._convert_with_markitdown(
            raw, filename, mime, conversation_id, user_id)
        if converted:
            return converted
        try:
            if mime == "application/pdf" or ext == "pdf":
                return self._extract_pdf_text(raw).strip()
            return self._convert_document_to_text(raw, filename, mime).strip()
        except Exception:
            logging.getLogger(__name__).debug(
                "Attachment text extraction failed for %s", filename, exc_info=True)
            return ""

    def _convert_with_markitdown(self, raw: bytes, filename: str, mime: str,
                                 conversation_id: str = "",
                                 user_id: str = "") -> str:
        try:
            import io
            from markitdown import MarkItDown
        except ImportError:
            return ""
        llm_client = None
        llm_model = ""
        ocr_service = str(
            self.config.get("attachment_ocr_llm_service")
            or os.getenv("PAWFLOW_MARKITDOWN_OCR_LLM_SERVICE", "")
            or "").strip()
        if ocr_service:
            llm_client, llm_model = self._markitdown_llm_adapter(
                ocr_service, user_id, conversation_id)
        try:
            kwargs = {"enable_plugins": bool(llm_client)}
            if llm_client:
                kwargs["llm_client"] = llm_client
                if llm_model:
                    kwargs["llm_model"] = llm_model
            md = MarkItDown(**kwargs)
            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
            attempts = [
                {"file_extension": f".{ext}"} if ext else {},
                {"filename": filename} if filename else {},
                {},
            ]
            last_error = None
            for call_kwargs in attempts:
                stream = io.BytesIO(raw)
                try:
                    result = md.convert_stream(stream, **call_kwargs)
                    text = (getattr(result, "text_content", "")
                            or getattr(result, "markdown", "")
                            or str(result or ""))
                    if text and text.strip():
                        return text.strip()
                except TypeError as exc:
                    last_error = exc
                    continue
            if last_error:
                raise last_error
        except Exception:
            logging.getLogger(__name__).debug(
                "MarkItDown conversion failed for %s", filename, exc_info=True)
        return ""

    def _markitdown_llm_adapter(self, service_id: str, user_id: str = "",
                                conversation_id: str = ""):
        try:
            from core.service_registry import ServiceRegistry
            reg = ServiceRegistry.get_instance()
            svc = reg.resolve(service_id, user_id=user_id, conv_id=conversation_id)
            if not svc or getattr(svc, "TYPE", "") != "llmConnection":
                return None, ""
            client = svc.get_client() if hasattr(svc, "get_client") else None
            if not client or not getattr(client, "supports_vision", False):
                return None, ""
            return (_PawFlowMarkItDownVisionClient(
                svc, user_id=user_id, conversation_id=conversation_id),
                getattr(svc, "default_model", "") or "pawflow-vision")
        except Exception:
            logging.getLogger(__name__).debug(
                "MarkItDown OCR LLM service resolution failed", exc_info=True)
            return None, ""


    @staticmethod
    def _convert_document_to_text(raw: bytes, filename: str, mime: str) -> str:
        """Convert office documents to text. Tries multiple libraries."""
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        # DOCX
        if ext == "docx" or "wordprocessingml" in mime:
            try:
                import io
                from docx import Document
                doc = Document(io.BytesIO(raw))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                if paragraphs:
                    return "\n\n".join(paragraphs)
            except ImportError:
                pass
            # Fallback: extract from zip XML
            try:
                import zipfile
                import io
                import re
                with zipfile.ZipFile(io.BytesIO(raw)) as z:
                    xml = z.read("word/document.xml").decode("utf-8")
                    text = re.sub(r'<[^>]+>', '', xml)
                    text = re.sub(r'\s+', ' ', text).strip()
                    if text:
                        return text
            except Exception:
                logging.getLogger(__name__).debug("Ignored exception", exc_info=True)
            raise ValueError("python-docx not available and XML extraction failed")

        # ODT
        if ext == "odt" or "opendocument.text" in mime:
            try:
                import zipfile
                import io
                import re
                with zipfile.ZipFile(io.BytesIO(raw)) as z:
                    xml = z.read("content.xml").decode("utf-8")
                    # Extract text between tags
                    text = re.sub(r'<[^>]+>', '\n', xml)
                    text = re.sub(r'\n{3,}', '\n\n', text).strip()
                    if text:
                        return text
            except Exception:
                logging.getLogger(__name__).debug("Ignored exception", exc_info=True)
            raise ValueError("ODT extraction failed")

        # XLSX
        if ext in ("xlsx", "xls") or "spreadsheet" in mime:
            try:
                import io
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
                sheets = []
                for ws in wb.worksheets:
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(cells):
                            rows.append("\t".join(cells))
                    if rows:
                        sheets.append(f"## Sheet: {ws.title}\n" + "\n".join(rows))
                wb.close()
                if sheets:
                    return "\n\n".join(sheets)
            except ImportError:
                pass
            raise ValueError("openpyxl not available")

        # PPTX
        if ext == "pptx" or "presentationml" in mime:
            try:
                import io
                from pptx import Presentation
                prs = Presentation(io.BytesIO(raw))
                slides = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    texts.append(t)
                    if texts:
                        slides.append(f"## Slide {i}\n" + "\n".join(texts))
                if slides:
                    return "\n\n".join(slides)
            except ImportError:
                pass
            raise ValueError("python-pptx not available")

        # ODS
        if ext == "ods" or "opendocument.spreadsheet" in mime:
            try:
                import zipfile
                import io
                import re
                with zipfile.ZipFile(io.BytesIO(raw)) as z:
                    xml = z.read("content.xml").decode("utf-8")
                    text = re.sub(r'<[^>]+>', '\t', xml)
                    text = re.sub(r'\t{3,}', '\n', text).strip()
                    if text:
                        return text
            except Exception:
                logging.getLogger(__name__).debug("Ignored exception", exc_info=True)
            raise ValueError("ODS extraction failed")

        # RTF
        if ext == "rtf" or "rtf" in mime:
            try:
                from striprtf.striprtf import rtf_to_text
                return rtf_to_text(raw.decode("utf-8", errors="replace"))
            except ImportError:
                # Basic RTF strip
                import re
                text = raw.decode("utf-8", errors="replace")
                text = re.sub(r'\\[a-z]+\d*\s?', '', text)
                text = re.sub(r'[{}]', '', text)
                return text.strip() or "(empty RTF)"

        # EPUB
        if ext == "epub" or "epub" in mime:
            try:
                import zipfile
                import io
                import re
                with zipfile.ZipFile(io.BytesIO(raw)) as z:
                    html_parts = []
                    for name in z.namelist():
                        if name.endswith((".html", ".xhtml", ".htm")):
                            html = z.read(name).decode("utf-8", errors="replace")
                            text = re.sub(r'<[^>]+>', ' ', html)
                            text = re.sub(r'\s+', ' ', text).strip()
                            if text:
                                html_parts.append(text)
                    if html_parts:
                        return "\n\n".join(html_parts)
            except Exception:
                logging.getLogger(__name__).debug("Ignored exception", exc_info=True)
            raise ValueError("EPUB extraction failed")

        raise ValueError(f"No converter for {ext}/{mime}")

    @staticmethod
    def _extract_pdf_text(raw_bytes: bytes) -> str:
        """Extract text from PDF bytes using available libraries."""
        # Try PyPDF2 first (most common)
        try:
            import io
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(raw_bytes))
            pages = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
            if pages:
                return "\n\n---\n\n".join(pages)
        except ImportError:
            pass
        except Exception:
            logging.getLogger(__name__).debug("Ignored exception", exc_info=True)

        # Try pdfminer
        try:
            import io
            from pdfminer.high_level import extract_text as _pdfminer_extract
            return _pdfminer_extract(io.BytesIO(raw_bytes))
        except ImportError:
            pass

        # Fallback: raw text extraction (basic)
        text = raw_bytes.decode("latin-1", errors="replace")
        # Extract readable strings (crude but works for simple PDFs)
        import re
        strings = re.findall(r'[\x20-\x7E]{10,}', text)
        if strings:
            return "\n".join(strings[:200])

        raise RuntimeError("No PDF library available (install PyPDF2 or pdfminer.six)")
